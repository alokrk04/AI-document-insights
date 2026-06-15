"""Parser for PPT and PPTX files using python-pptx."""
from pathlib import Path


def parse_ppt(file_path: Path) -> dict:
    """Extract text from PowerPoint preserving slide structure."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    sections = []
    full_text = ""

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)

        if slide_texts:
            slide_content = "\n".join(slide_texts)
            sections.append({
                "page_number": slide_num,
                "section": f"Slide {slide_num}",
                "text": slide_content,
                "source_type": "ppt",
            })
            full_text += f"--- Slide {slide_num} ---\n{slide_content}\n\n"

    return {
        "pages": sections,
        "full_text": full_text.strip(),
        "page_count": len(sections),
        "source_type": "ppt",
        "summary": f"PowerPoint with {len(sections)} slides.",
    }
